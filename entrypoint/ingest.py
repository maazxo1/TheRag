"""
Document ingestion pipeline.

Steps:
  1. Load all files from ./docs/ (PDF, TXT, DOCX, MD, ...)
  2. Build parent + child nodes via small-to-big chunking
  3. Embed child nodes → store in ChromaDB (persistent vector index)
  4. Build BM25 keyword index from child nodes
  5. Pickle BM25 index + node maps to ./bm25_index/
"""

import os
import pickle
import time
import tempfile
import shutil

import chromadb
import docx2txt
import fitz  # PyMuPDF
import pptx
from llama_index.core import SimpleDirectoryReader, VectorStoreIndex, StorageContext
from llama_index.core.readers.base import BaseReader
from llama_index.core.schema import Document
from llama_index.embeddings.ollama import OllamaEmbedding
from llama_index.vector_stores.chroma import ChromaVectorStore
from rank_bm25 import BM25Okapi

import config
from src.pipelines.chunking_pipeline import build_small_to_big_nodes


def _is_clean_text(text: str, min_printable_ratio: float = 0.85) -> bool:
    """Return False if the chunk looks like binary/PDF garbage rather than real text."""
    stripped = text.strip()
    if len(stripped) < 20:
        return False
    # Use str.isprintable() alone — covers all valid Unicode including accented
    # characters, CJK, Arabic, em-dashes, curly quotes, etc. The previous
    # `ord(c) < 128` restriction incorrectly rejected legitimate non-ASCII text.
    printable = sum(1 for c in stripped if c.isprintable())
    return (printable / len(stripped)) >= min_printable_ratio


class PdfReader(BaseReader):
    """Reads PDFs via PyMuPDF — handles complex layouts, diagrams, and fonts."""
    def load_data(self, file, extra_info=None):
        doc = fitz.open(str(file))
        pages = []
        for page in doc:
            text = page.get_text()
            if text.strip():
                pages.append(text)
        doc.close()
        text = "\n\n".join(pages)
        metadata = extra_info or {}
        metadata["file_name"] = os.path.basename(str(file))
        return [Document(text=text, metadata=metadata)]


class DocxReader(BaseReader):
    """Reads .docx files via docx2txt — no llama-index-readers-file required."""
    def load_data(self, file, extra_info=None):
        text = docx2txt.process(str(file))
        metadata = extra_info or {}
        metadata["file_name"] = os.path.basename(str(file))
        return [Document(text=text, metadata=metadata)]


class PptxReader(BaseReader):
    """
    Reads .pptx files via python-pptx.

    Extracts text from every shape on every slide — titles, text boxes,
    table cells, and notes. One Document per file (not per slide) so
    the chunker can create coherent parent/child chunks across slides.
    """
    def load_data(self, file, extra_info=None):
        presentation = pptx.Presentation(str(file))
        slide_texts = []
        for slide_num, slide in enumerate(presentation.slides, start=1):
            parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            parts.append(line)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            parts.append(" | ".join(cells))
            # Include speaker notes if present
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    parts.append(f"[Notes] {notes_text}")
            if parts:
                slide_texts.append(f"[Slide {slide_num}]\n" + "\n".join(parts))

        text = "\n\n".join(slide_texts)
        metadata = extra_info or {}
        metadata["file_name"] = os.path.basename(str(file))
        return [Document(text=text, metadata=metadata)]

# Paths for persisted artefacts
_BM25_PATH = os.path.join(config.BM25_INDEX_PATH, "bm25_index.pkl")
_PARENT_PATH = os.path.join(config.BM25_INDEX_PATH, "parent_nodes.pkl")
_CHILD_PATH = os.path.join(config.BM25_INDEX_PATH, "child_nodes.pkl")


def ingest_documents(data_dir: str = config.DATA_DIR, force: bool = False):
    """
    Ingest all supported documents from data_dir.

    Args:
        data_dir: folder containing user documents
        force:    re-ingest even if an existing index is detected

    Returns:
        index           — LlamaIndex VectorStoreIndex (over child chunks)
        bm25            — BM25Okapi instance
        child_nodes     — list of TextNode (child chunks)
        parent_node_map — {node_id: TextNode} for parent chunks
    """
    # ── 1. Load documents ───────────────────────────────────────────────────────
    print(f"\n[Ingest] Loading documents from '{data_dir}' ...")
    reader = SimpleDirectoryReader(
        input_dir=data_dir,
        recursive=True,
        filename_as_id=True,
        file_extractor={
            ".pdf": PdfReader(),
            ".docx": DocxReader(),
            ".pptx": PptxReader(),
        },
    )
    documents = reader.load_data()
    if not documents:
        raise FileNotFoundError(
            f"No supported documents found in '{data_dir}'. "
            "Drop PDFs or text files there and re-run."
        )
    print(f"[Ingest] Loaded {len(documents)} document(s).")

    # ── 2. Small-to-Big chunking ────────────────────────────────────────────────
    print("[Ingest] Chunking documents (small-to-big) ...")
    t0 = time.time()
    parent_nodes, child_nodes, parent_node_map = build_small_to_big_nodes(documents)

    # Filter out binary/garbage chunks (common with complex PDFs).
    # Parents are filtered first; then children are filtered to match so that
    # no child node is left pointing to a parent that no longer exists (orphan).
    before_children = len(child_nodes)
    before_parents = len(parent_node_map)

    parent_node_map = {
        pid: node for pid, node in parent_node_map.items()
        if _is_clean_text(node.get_content())
    }
    # Keep only children whose content is clean AND whose parent survived filtering.
    child_nodes = [
        n for n in child_nodes
        if _is_clean_text(n.get_content())
        and n.metadata.get("parent_id", n.node_id) in parent_node_map
    ]

    dropped_c = before_children - len(child_nodes)
    dropped_p = before_parents - len(parent_node_map)
    if dropped_c or dropped_p:
        print(
            f"[Ingest] Dropped {dropped_p} parent chunk(s) and {dropped_c} child chunk(s) "
            "(binary/non-text or orphaned content)."
        )

    print(
        f"[Ingest] {len(parent_node_map)} parent chunks  |  "
        f"{len(child_nodes)} child chunks  |  {time.time()-t0:.1f}s"
    )

    if not child_nodes:
        raise ValueError(
            "No usable text chunks produced — the document may be too short, "
            "empty, or contain only non-text content."
        )

    # ── 3. ChromaDB vector index ────────────────────────────────────────────────
    print("[Ingest] Building vector index (embedding child chunks) ...")
    os.makedirs(config.DB_PATH, exist_ok=True)
    chroma_client = chromadb.PersistentClient(path=config.DB_PATH)

    # Drop old collection on force re-ingest so we don't accumulate stale data
    if force:
        try:
            chroma_client.delete_collection(config.CHROMA_COLLECTION)
        except Exception:
            pass

    chroma_collection = chroma_client.get_or_create_collection(config.CHROMA_COLLECTION)
    embed_model = OllamaEmbedding(
        model_name=config.EMBEDDING_MODEL,
        base_url=config.OLLAMA_URL,
    )
    vector_store = ChromaVectorStore(chroma_collection=chroma_collection)
    storage_context = StorageContext.from_defaults(vector_store=vector_store)

    t0 = time.time()
    index = VectorStoreIndex(
        child_nodes,
        storage_context=storage_context,
        embed_model=embed_model,
        show_progress=True,
    )
    print(f"[Ingest] Vector index ready  |  {time.time()-t0:.1f}s")

    # ── 4. BM25 keyword index ───────────────────────────────────────────────────
    print("[Ingest] Building BM25 keyword index ...")
    corpus = [node.get_content().lower().split() for node in child_nodes]
    bm25 = BM25Okapi(corpus)

    # ── 5. Persist BM25 artefacts (atomic write → rename so a crash mid-write
    #       never leaves BM25 and ChromaDB in an inconsistent state) ─────────────
    os.makedirs(config.BM25_INDEX_PATH, exist_ok=True)
    for path, obj in [(_BM25_PATH, bm25), (_PARENT_PATH, parent_node_map), (_CHILD_PATH, child_nodes)]:
        tmp = path + ".tmp"
        with open(tmp, "wb") as f:
            pickle.dump(obj, f)
        os.replace(tmp, path)  # atomic on POSIX; best-effort on Windows
    print(f"[Ingest] BM25 index + node maps saved to '{config.BM25_INDEX_PATH}'")

    print("[Ingest] Done.\n")
    return index, bm25, child_nodes, parent_node_map


def load_existing_index():
    """
    Load a previously ingested index from disk without re-embedding.
    Raises FileNotFoundError if no index exists yet.
    """
    for path in (_BM25_PATH, _PARENT_PATH, _CHILD_PATH):
        if not os.path.exists(path):
            raise FileNotFoundError(
                f"No persisted index found at '{path}'. Run ingest_documents() first."
            )

    chroma_client = chromadb.PersistentClient(path=config.DB_PATH)
    chroma_collection = chroma_client.get_or_create_collection(config.CHROMA_COLLECTION)
    embed_model = OllamaEmbedding(
        model_name=config.EMBEDDING_MODEL,
        base_url=config.OLLAMA_URL,
    )
    vector_store = ChromaVectorStore(chroma_collection=chroma_collection)
    index = VectorStoreIndex.from_vector_store(
        vector_store,
        embed_model=embed_model,
    )

    with open(_BM25_PATH, "rb") as f:
        bm25 = pickle.load(f)
    with open(_PARENT_PATH, "rb") as f:
        parent_node_map = pickle.load(f)
    with open(_CHILD_PATH, "rb") as f:
        child_nodes = pickle.load(f)

    print(
        f"[Ingest] Loaded existing index — "
        f"{len(child_nodes)} child chunks, {len(parent_node_map)} parent chunks."
    )
    return index, bm25, child_nodes, parent_node_map


def _clear_data_dir() -> None:
    """Remove only files (not subdirectories) from DATA_DIR before re-ingest."""
    os.makedirs(config.DATA_DIR, exist_ok=True)
    for entry in os.listdir(config.DATA_DIR):
        full = os.path.join(config.DATA_DIR, entry)
        if os.path.isfile(full):
            os.remove(full)


def ingest_from_file(file_path: str):
    """
    Ingest a single uploaded file (PDF / DOCX / TXT).
    Clears docs/ first so uploaded file is the ONLY source — no mixing.
    """
    _clear_data_dir()
    dest = os.path.join(config.DATA_DIR, os.path.basename(file_path))
    shutil.copy2(file_path, dest)
    return ingest_documents(force=True)


def ingest_from_text(text: str, filename: str = "pasted_document.txt"):
    """
    Ingest raw pasted text by writing it to a temp file in DATA_DIR.
    Clears docs/ first so pasted text is the ONLY source.
    """
    _clear_data_dir()
    dest = os.path.join(config.DATA_DIR, filename)
    with open(dest, "w", encoding="utf-8") as f:
        f.write(text)
    return ingest_documents(force=True)


if __name__ == "__main__":
    ingest_documents(force=True)
